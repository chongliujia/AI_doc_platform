from typing import Dict, Any, Optional, List
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.enum.shapes import MSO_SHAPE
import logging
import re

from .ai_service_factory import AIServiceFactory
from .image_service import ImageService

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTGenerator:
    def __init__(self, ai_service_type: str = "deepseek"):
        self.templates_dir = os.path.join(os.path.dirname(__file__), "../templates/ppt_templates")
        # 修正路径：生成文档目录和app是同级的
        self.output_dir = os.path.abspath(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "generated_docs"))
        logger.info(f"PPT生成器输出目录: {self.output_dir}")
        os.makedirs(self.output_dir, exist_ok=True)
        
        # 确保目录权限正确
        try:
            os.chmod(self.output_dir, 0o777)
        except Exception as e:
            logger.warning(f"无法修改目录权限: {e}")
        
        # 定义配色方案
        self.COLORS = {
            'primary': RGBColor(41, 128, 185),    # 主色调：蓝色
            'secondary': RGBColor(44, 62, 80),    # 次要色：深灰
            'accent': RGBColor(231, 76, 60),      # 强调色：红色
            'light': RGBColor(236, 240, 241),     # 浅色：近白
            'dark': RGBColor(52, 73, 94)          # 深色：深灰蓝
        }
        
        self.prs = None
        self.ai_service = AIServiceFactory.create_service(ai_service_type)
        self.image_service = ImageService()

    def _sanitize_filename(self, name: str) -> str:
        """
        将主题名称转换为有效的文件名
        
        Args:
            name: 原始主题名称
            
        Returns:
            有效的文件名
        """
        # 移除不合法的文件名字符
        # 移除特殊字符，保留字母、数字、空格和一些基本标点
        sanitized = re.sub(r'[\\/*?:"<>|]', '', name)
        # 移除换行符并替换多个空格为单个空格
        sanitized = re.sub(r'\s+', ' ', sanitized.replace('\n', ' '))
        # 限制长度
        max_length = 100
        if len(sanitized) > max_length:
            # 如果太长，截断并添加提示
            sanitized = sanitized[:max_length] + "..."
        # 确保文件名不为空
        if not sanitized.strip():
            sanitized = "AI_presentation"
        return sanitized.strip()

    def generate(self, topic: str, outline: List[Dict[str, Any]], template_id: Optional[str] = None, max_pages: Optional[int] = None) -> Optional[str]:
        """
        生成PPT文档
        
        Args:
            topic: 主题
            outline: 大纲内容
            template_id: 模板ID（现在不使用，保留参数以保持接口兼容）
            max_pages: 限制生成的最大页数
        """
        try:
            logger.info(f"开始生成PPT: 主题='{topic}', 章节数={len(outline)}")
            if template_id:
                logger.info(f"使用模板: {template_id}")
            if max_pages:
                logger.info(f"应用页数限制: 最大 {max_pages} 张幻灯片")
            
            # 创建新的演示文稿
            self.prs = Presentation()
            self.prs.slide_width = Inches(16)
            self.prs.slide_height = Inches(9)
            logger.info(f"创建新的演示文稿: 宽度={self.prs.slide_width}, 高度={self.prs.slide_height}")
            
            # 计算固定幻灯片数量（标题、目录和结束幻灯片）
            fixed_slides = 3
            
            # 如果有页数限制，计算可用于章节内容的幻灯片数量
            remaining_slides = None
            if max_pages:
                if max_pages <= fixed_slides:
                    logger.warning(f"页数限制({max_pages})太小，至少需要{fixed_slides}张基本幻灯片，将使用最小值")
                    max_pages = fixed_slides + 1
                remaining_slides = max_pages - fixed_slides
            
            # 添加标题幻灯片
            logger.info("添加标题幻灯片")
            self._add_title_slide(topic)
            
            # 添加目录幻灯片
            logger.info("添加目录幻灯片")
            self._add_toc_slide(outline)
            
            # 计算章节标题幻灯片数量
            section_title_slides = len(outline)
            content_slides_count = 0
            
            # 如果有页数限制且章节标题幻灯片已经超过限制，裁剪章节
            if max_pages and (section_title_slides + fixed_slides > max_pages):
                # 我们需要至少保留一个章节
                available_sections = max_pages - fixed_slides
                if available_sections < 1:
                    available_sections = 1
                
                logger.warning(f"章节数({len(outline)})超过可用页数限制，将裁剪至{available_sections}个章节")
                outline = outline[:available_sections]
                section_title_slides = len(outline)
            
            # 如果有页数限制，计算每个章节内容可用的幻灯片数量
            available_content_slides = None
            if max_pages:
                available_content_slides = max_pages - fixed_slides - section_title_slides
                if available_content_slides < len(outline):
                    available_content_slides = len(outline)  # 确保每个章节至少有一张内容幻灯片
                
                logger.info(f"可用于内容的幻灯片数量: {available_content_slides} (总限制: {max_pages}, 已用: {fixed_slides + section_title_slides})")
                
                # 计算每个章节平均可用的幻灯片数量
                slides_per_section = available_content_slides // len(outline)
                remaining = available_content_slides % len(outline)
                
                # 分配给每个章节的幻灯片数量
                section_slide_allocation = [slides_per_section] * len(outline)
                # 将剩余的幻灯片分配给前面的章节
                for i in range(remaining):
                    section_slide_allocation[i] += 1
                
                logger.info(f"章节幻灯片分配: {section_slide_allocation}")
            
            # 处理每个章节
            total_slides = 2  # 已添加标题和目录幻灯片
            for section_index, section in enumerate(outline):
                section_title = section["title"]
                logger.info(f"处理章节 {section_index+1}/{len(outline)}: '{section_title}'")
                
                # 添加章节标题幻灯片
                self._add_section_title_slide(section_title)
                total_slides += 1
                
                # 添加章节内容幻灯片
                slides = section.get("slides", [])
                
                # 如果有页数限制，限制本章节的幻灯片数量
                if max_pages and section_slide_allocation:
                    allowed_slides = section_slide_allocation[section_index]
                    if len(slides) > allowed_slides:
                        logger.info(f"章节 '{section_title}' 幻灯片数量({len(slides)})超过分配值({allowed_slides})，将裁剪")
                        slides = slides[:allowed_slides]
                        
                logger.info(f"  章节 '{section_title}' 包含 {len(slides)} 张幻灯片")
                
                for slide_index, slide_content in enumerate(slides):
                    # 检查是否已达到总页数限制
                    if max_pages and total_slides + 1 >= max_pages:  # +1 考虑最后的结束幻灯片
                        logger.warning(f"已达到页数限制({max_pages})，停止添加更多幻灯片")
                        break
                        
                    slide_title = slide_content.get("title", "未知标题")
                    slide_type = slide_content.get("type", "content")
                    logger.info(f"  添加幻灯片 {slide_index+1}/{len(slides)}: '{slide_title}' (类型: {slide_type})")
                    self._add_content_slide(slide_content, topic, section_title)
                    total_slides += 1
                    content_slides_count += 1
            
            # 添加结束幻灯片
            logger.info("添加结束幻灯片")
            self._add_ending_slide(topic)
            total_slides += 1
            
            # 保存文件 - 使用安全的文件名
            sanitized_topic = self._sanitize_filename(topic)
            # 确保文件保存到 downloads 子目录
            downloads_dir = os.path.join(self.output_dir, "downloads")
            os.makedirs(downloads_dir, exist_ok=True)
            output_path = os.path.join(downloads_dir, f"{sanitized_topic}_presentation.pptx")
            logger.info(f"保存PPT文件: {output_path}")
            self.prs.save(output_path)
            
            logger.info(f"PPT生成完成: 共 {total_slides} 张幻灯片 (包含内容幻灯片 {content_slides_count} 张), 保存至: {output_path}")
            if max_pages:
                logger.info(f"页数限制: {max_pages} 张, 实际生成: {total_slides} 张")
            
            # 清理临时文件
            self.image_service.cleanup()
            
            return output_path
            
        except Exception as e:
            logger.error(f"生成PPT时出错: {str(e)}")
            return None
        finally:
            self.prs = None

    def _add_points(self, text_frame: TextFrame, points: List[Dict[str, Any]]) -> None:
        """添加要点和详细说明，自动调整字体大小和布局"""
        # 要点数量的阈值，超过此数量则减小字体大小
        points_threshold = 3
        details_threshold = 2
        
        # 根据要点数量调整字体大小
        main_point_size = Pt(28)
        detail_point_size = Pt(20)
        
        if len(points) > points_threshold:
            # 要点较多，减小字体
            main_point_size = Pt(24)
            detail_point_size = Pt(18)
            
        for i, point in enumerate(points):
            # 添加主要要点
            p = text_frame.add_paragraph()
            p.text = "▪ " + point.get("main", "")  # 添加项目符号
            p.level = 0
            p.font.name = '微软雅黑'
            p.font.size = main_point_size
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            p.space_after = Pt(8) if len(points) > points_threshold else Pt(12)  # 根据内容量调整间距
            
            # 获取详细说明
            details = point.get("details", [])
            
            # 如果详细说明太多，减小字体大小和间距
            detail_font_size = detail_point_size
            if len(details) > details_threshold:
                detail_font_size = Pt(16)
            
            # 添加详细说明
            for detail in details:
                p = text_frame.add_paragraph()
                p.text = "• " + detail
                p.level = 1
                p.font.name = '微软雅黑'
                p.font.size = detail_font_size
                p.font.color.rgb = self.COLORS['secondary']
                p.space_before = Pt(4) if len(details) > details_threshold else Pt(6)  # 根据内容量调整间距
                p.space_after = Pt(4) if len(details) > details_threshold else Pt(6)   # 根据内容量调整间距

    def _add_content_slide(self, content: Dict[str, Any], topic: str, section_title: str) -> None:
        """添加内容幻灯片，根据内容量自动调整布局"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 添加标题和装饰线条
        title = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(14), Inches(1)
        )
        title.text_frame.text = content.get("title", "")
        p = title.text_frame.paragraphs[0]
        p.font.name = '微软雅黑'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['dark']
        p.alignment = PP_ALIGN.LEFT
        
        # 添加装饰线条（使用矩形作为线条）
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.3), Inches(14), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['primary']
        
        # 获取幻灯片内容
        slide_title = content.get("title", "")
        slide_type = content.get("type", "content")
        
        # 使用AI服务生成幻灯片内容
        slide_content = self.ai_service.generate_slide_content(topic, section_title, slide_title, slide_type)
        
        # 获取要点数量，判断是否需要双列布局
        points = slide_content.get("points", [])
        total_points = len(points)
        total_details = sum(len(point.get("details", [])) for point in points)
        
        # 自动判断布局方式
        use_two_column = total_points > 3 or total_details > 6
        
        if slide_type == "two_column" or use_two_column:
            # 如果指定双列或内容较多，使用两栏布局
            if slide_type == "two_column" and "left_points" in slide_content and "right_points" in slide_content:
                # 使用预定义的左右栏内容
                left_content = slide.shapes.add_textbox(
                    Inches(1), Inches(1.5), Inches(6.5), Inches(6)
                )
                right_content = slide.shapes.add_textbox(
                    Inches(8), Inches(1.5), Inches(6.5), Inches(6)
                )
                self._add_points(left_content.text_frame, slide_content.get("left_points", []))
                self._add_points(right_content.text_frame, slide_content.get("right_points", []))
            else:
                # 自动将内容分为左右两栏
                left_points = points[:len(points)//2 + len(points)%2]  # 左侧放置一半+余数
                right_points = points[len(points)//2 + len(points)%2:]  # 右侧放置剩余部分
                
                left_content = slide.shapes.add_textbox(
                    Inches(1), Inches(1.5), Inches(6.5), Inches(6)
                )
                right_content = slide.shapes.add_textbox(
                    Inches(8), Inches(1.5), Inches(6.5), Inches(6)
                )
                self._add_points(left_content.text_frame, left_points)
                self._add_points(right_content.text_frame, right_points)
        else:
            # 普通单列布局
            text_content = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), 
                Inches(14 if slide_type == "content" else 8), 
                Inches(6.5)  # 增加高度
            )
            self._add_points(text_content.text_frame, points)
            
            if slide_type == "image_content":
                # 获取图片描述
                image_desc = slide_content.get("image_description", "")
                if image_desc:
                    # 使用图片服务获取图片
                    image_path = self.image_service.get_image_for_slide(image_desc)
                    
                    if image_path and os.path.exists(image_path):
                        try:
                            # 添加图片到幻灯片
                            slide.shapes.add_picture(
                                image_path,
                                Inches(9.5), Inches(2),  # 位置
                                width=Inches(5.5)        # 宽度，高度自动按比例调整
                            )
                            logger.info(f"已添加图片到幻灯片 '{slide_title}'")
                        except Exception as e:
                            logger.error(f"添加图片时出错: {str(e)}")
                            # 添加图片说明文本框
                            image_note = slide.shapes.add_textbox(
                                Inches(9.5), Inches(3.5), Inches(5.5), Inches(1)
                            )
                            image_note.text_frame.text = f"[图片: {image_desc}]"
                            p = image_note.text_frame.paragraphs[0]
                            p.font.italic = True
                            p.font.color.rgb = self.COLORS['secondary']
                    else:
                        # 如果无法获取图片，添加图片说明文本框
                        image_note = slide.shapes.add_textbox(
                            Inches(9.5), Inches(3.5), Inches(5.5), Inches(1)
                        )
                        image_note.text_frame.text = f"[图片: {image_desc}]"
                        p = image_note.text_frame.paragraphs[0]
                        p.font.italic = True
                        p.font.color.rgb = self.COLORS['secondary']
        
        # 添加注释
        if notes := content.get("notes"):
            slide.notes_slide.notes_text_frame.text = notes

    def _add_title_slide(self, topic: str) -> None:
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 使用标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic
        subtitle.text = "AI自动生成的演示文稿"
        
        # 设置标题样式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['primary']
        
        # 设置副标题样式
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.italic = True
                run.font.color.rgb = self.COLORS['secondary']

    def _add_toc_slide(self, outline: List[Dict[str, Any]]) -> None:
        """添加目录幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 使用空白布局
        
        # 添加标题
        title = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(14), Inches(1)
        )
        title.text_frame.text = "目录"
        title_p = title.text_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = self.COLORS['primary']
        
        # 添加装饰线条
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # 修改这里
            Inches(2), Inches(1.8), Inches(12), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['primary']
        
        # 创建两列布局
        left_content = slide.shapes.add_textbox(
            Inches(2), Inches(2.2), Inches(5.5), Inches(5)
        )
        right_content = slide.shapes.add_textbox(
            Inches(8.5), Inches(2.2), Inches(5.5), Inches(5)
        )
        
        # 分配目录项到两列
        mid_point = len(outline) // 2 + len(outline) % 2
        
        # 添加左列目录项
        tf_left = left_content.text_frame
        for i, section in enumerate(outline[:mid_point], 1):
            p = tf_left.add_paragraph()
            p.text = f"{i}. {section['title']}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['secondary']
            p.space_after = Pt(20)  # 增加段落间距
        
        # 添加右列目录项
        tf_right = right_content.text_frame
        for i, section in enumerate(outline[mid_point:], mid_point + 1):
            p = tf_right.add_paragraph()
            p.text = f"{i}. {section['title']}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['secondary']
            p.space_after = Pt(20)  # 增加段落间距

    def _add_section_title_slide(self, section_title: str) -> None:
        """添加章节标题幻灯片"""
        slide_layout = self.prs.slide_layouts[2]  # 使用章节标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = section_title
        
        # 设置标题样式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['primary']

    def _add_ending_slide(self, topic: str) -> None:
        """添加结束幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "总结与问答"
        
        # 设置标题样式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['primary']
        
        # 添加总结内容
        tf = content.text_frame
        
        summary_points = [
            f"我们探讨了{topic}的多个关键方面",
            f"理解这些概念对把握{topic}的本质至关重要",
            f"未来发展将带来更多机遇和挑战",
            f"感谢您的关注！有任何问题欢迎提问"
        ]
        
        for point in summary_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
            # 设置段落样式
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.color.rgb = self.COLORS['secondary'] 